"""Native PowerPoint parser backend.

The backend extracts presentation text, native tables, and picture shapes
without rendering slides or calling an OCR/model provider.  Its output matches
the semantic ``extraction`` contract used by the document parsers so existing
validation, image filtering, indexing, and artifact writers can consume it.
"""

from __future__ import annotations

from dataclasses import dataclass
import hashlib
from pathlib import Path
import re
from typing import Any, Iterable

from ....models import DataObject, ParsedData, ParsedTable, ParseResult
from ....utils.paths import portable_path
from ..contracts import AXIOM_NATIVE_BLOCK_SOURCE
from ._tabular import normalize_headers


PPTX_EXTENSIONS = frozenset({".pptx"})
_GENERIC_PICTURE_NAME = re.compile(r"^(?:picture|image|graphic)\s+\d+$", re.I)
_UNSUPPORTED_CONTENT_SHAPES = {
    "CHART",
    "DIAGRAM",
    "EMBEDDED_OLE_OBJECT",
    "LINKED_OLE_OBJECT",
    "MEDIA",
}


@dataclass(frozen=True)
class PptxConfig:
    """Asset settings for native PowerPoint extraction."""

    extract_images: bool = True
    output_dir: str | None = None
    project_root: str | None = None

    @classmethod
    def from_mapping(cls, config: dict[str, Any] | None) -> "PptxConfig":
        config = config or {}
        return cls(
            extract_images=bool(config.get("extract_images", True)),
            output_dir=_optional_str(config.get("output_dir")),
            project_root=_optional_str(config.get("project_root")),
        )


@dataclass
class _PptxState:
    source_blocks: list[dict[str, Any]]
    text_parts: list[str]
    text_citations: list[str]
    tables: list[dict[str, Any]]
    parsed_tables: list[ParsedTable]
    figures: list[dict[str, Any]]
    image_files: list[dict[str, Any]]
    unsupported_shape_types: dict[str, int]
    first_title: str | None = None


class PptxParserBackend:
    """Parse PPTX files through their native Open XML structure."""

    supported_extensions = PPTX_EXTENSIONS
    backend_name = "pptx"

    def __init__(self, config: PptxConfig | None = None) -> None:
        self.config = config or PptxConfig()

    def parse(self, path: str | Path, data_object: DataObject) -> ParseResult:
        file_path = Path(path)
        try:
            parsed = self._parse(file_path, data_object)
        except Exception as exc:
            return ParseResult.failed(
                data_object.object_id,
                self.backend_name,
                exc,
                route=self.backend_name,
            )
        return ParseResult.success(
            data_object.object_id,
            self.backend_name,
            parsed,
            route=self.backend_name,
        )

    def _parse(self, path: Path, data_object: DataObject) -> ParsedData:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise RuntimeError(
                "python-pptx is required to parse .pptx files"
            ) from exc

        presentation = Presentation(str(path))
        state = _PptxState([], [], [], [], [], [], [], {})
        saved_assets: dict[str, dict[str, Any]] = {}
        slide_width = int(presentation.slide_width)
        slide_height = int(presentation.slide_height)
        page_box = [0, 0, slide_width, slide_height]

        for slide_index, slide in enumerate(presentation.slides):
            for shape in _ordered_shapes(slide.shapes):
                self._parse_shape(
                    shape,
                    slide_index=slide_index,
                    page_box=page_box,
                    data_object=data_object,
                    state=state,
                    saved_assets=saved_assets,
                )

        main_text = "\n\n".join(state.text_parts)
        extraction = {
            "document_type": "presentation",
            "language": None,
            "title": state.first_title,
            "title_citations": _title_citations(state),
            "main_text": main_text,
            "main_text_citations": state.text_citations,
            "tables": state.tables,
            "figures": state.figures,
            "formulas": [],
            "formulas_citations": [],
        }
        unsupported_count = sum(state.unsupported_shape_types.values())
        return ParsedData(
            object_id=data_object.object_id,
            source_uri=data_object.uri,
            source_format="pptx",
            rows=[
                {
                    "extraction": extraction,
                    "text": main_text,
                    "source_blocks": state.source_blocks,
                    "reading_order": [
                        block["component_id"] for block in state.source_blocks
                    ],
                }
            ],
            text=main_text,
            tables=state.parsed_tables,
            metadata={
                "parser": self.backend_name,
                "page_count": len(presentation.slides),
                "source_block_count": len(state.source_blocks),
                "table_count": len(state.tables),
                "figure_count": len(state.figures),
                "formula_count": 0,
                "image_count": sum(
                    item.get("status") == "saved" for item in state.image_files
                ),
                "image_files": state.image_files,
                "unsupported_shape_count": unsupported_count,
                "unsupported_shape_types": dict(
                    sorted(state.unsupported_shape_types.items())
                ),
                "reading_order_source": "pptx_native",
                "reading_order_complete": unsupported_count == 0,
                "notes_included": False,
            },
        )

    def _parse_shape(
        self,
        shape: Any,
        *,
        slide_index: int,
        page_box: list[int],
        data_object: DataObject,
        state: _PptxState,
        saved_assets: dict[str, dict[str, Any]],
    ) -> None:
        shape_type = _enum_name(getattr(shape, "shape_type", "UNKNOWN"))
        if shape_type == "GROUP":
            for child in _ordered_shapes(shape.shapes):
                self._parse_shape(
                    child,
                    slide_index=slide_index,
                    page_box=page_box,
                    data_object=data_object,
                    state=state,
                    saved_assets=saved_assets,
                )
            return

        if bool(getattr(shape, "has_table", False)):
            self._parse_table(shape, slide_index, page_box, state)
            return

        if shape_type == "PICTURE":
            self._parse_picture(
                shape,
                slide_index,
                page_box,
                data_object,
                state,
                saved_assets,
            )
            return

        if bool(getattr(shape, "has_text_frame", False)):
            text = _text_frame_text(shape.text_frame)
            if text:
                block_type = "Title" if _is_title_shape(shape) else "Text"
                component_id = _append_block(
                    state,
                    slide_index=slide_index,
                    block_type=block_type,
                    text=text,
                    shape=shape,
                    page_box=page_box,
                )
                state.text_parts.append(text)
                state.text_citations.append(component_id)
                if (
                    slide_index == 0
                    and block_type == "Title"
                    and state.first_title is None
                ):
                    state.first_title = text.splitlines()[0].strip()
            return

        if bool(getattr(shape, "has_chart", False)):
            shape_type = "CHART"
        if shape_type in _UNSUPPORTED_CONTENT_SHAPES:
            state.unsupported_shape_types[shape_type] = (
                state.unsupported_shape_types.get(shape_type, 0) + 1
            )

    def _parse_table(
        self,
        shape: Any,
        slide_index: int,
        page_box: list[int],
        state: _PptxState,
    ) -> None:
        raw_rows = [
            [
                "" if bool(getattr(cell, "is_spanned", False)) else _clean_text(cell.text)
                for cell in row.cells
            ]
            for row in shape.table.rows
        ]
        width = max((len(row) for row in raw_rows), default=0)
        if width == 0:
            return
        raw_rows = [row + [""] * (width - len(row)) for row in raw_rows]
        headers = normalize_headers(raw_rows[0], width=width)
        rows = raw_rows[1:]
        content = _markdown_table(headers, rows)
        component_id = _append_block(
            state,
            slide_index=slide_index,
            block_type="Table",
            text=content,
            shape=shape,
            page_box=page_box,
        )
        name = f"Slide {slide_index + 1} Table {len(state.tables) + 1}"
        state.tables.append(
            {
                "caption": "",
                "caption_citations": [],
                "content": content,
                "content_citations": [component_id],
                "content_format": "markdown",
                "headers": headers,
                "rows": rows,
                "source_ref": component_id,
            }
        )
        state.parsed_tables.append(
            ParsedTable(
                name=name,
                source_ref=component_id,
                headers=headers,
                rows=rows,
                metadata={
                    "parser": self.backend_name,
                    "slide_index": slide_index,
                },
            )
        )

    def _parse_picture(
        self,
        shape: Any,
        slide_index: int,
        page_box: list[int],
        data_object: DataObject,
        state: _PptxState,
        saved_assets: dict[str, dict[str, Any]],
    ) -> None:
        description = _picture_description(shape)
        component_id = _append_block(
            state,
            slide_index=slide_index,
            block_type="Figure",
            text=description,
            shape=shape,
            page_box=page_box,
        )
        state.figures.append(
            {
                "caption": "",
                "caption_citations": [],
                "description": description,
                "description_citations": [component_id],
            }
        )
        state.image_files.append(
            self._image_asset(
                shape,
                slide_index=slide_index,
                component_id=component_id,
                data_object=data_object,
                saved_assets=saved_assets,
            )
        )

    def _image_asset(
        self,
        shape: Any,
        *,
        slide_index: int,
        component_id: str,
        data_object: DataObject,
        saved_assets: dict[str, dict[str, Any]],
    ) -> dict[str, Any]:
        image = shape.image
        blob = bytes(image.blob)
        digest = hashlib.sha256(blob).hexdigest()
        extension = _safe_extension(getattr(image, "ext", None))
        file_name = f"{digest}.{extension}"
        saved = saved_assets.get(digest)
        if saved is None:
            saved = self._save_image(data_object.object_id, file_name, blob)
            saved_assets[digest] = saved
        return {
            **saved,
            "page": slide_index,
            "source_ref": component_id,
            "sha256": digest,
        }

    def _save_image(
        self,
        object_id: str,
        file_name: str,
        blob: bytes,
    ) -> dict[str, Any]:
        if not self.config.extract_images or not self.config.output_dir:
            return {
                "name": file_name,
                "path": None,
                "status": "not_saved",
            }
        image_dir = Path(self.config.output_dir) / object_id / "images"
        image_dir.mkdir(parents=True, exist_ok=True)
        output_path = image_dir / file_name
        if not output_path.exists():
            temporary = output_path.with_name(f".{output_path.name}.tmp")
            temporary.write_bytes(blob)
            temporary.replace(output_path)
        return {
            "name": file_name,
            "path": portable_path(output_path, self.config.project_root),
            "status": "saved",
        }


def _ordered_shapes(shapes: Iterable[Any]) -> list[Any]:
    indexed = list(enumerate(shapes))
    return [
        shape
        for _, shape in sorted(
            indexed,
            key=lambda item: (
                _shape_priority(item[1]),
                _int_value(getattr(item[1], "top", 0)),
                _int_value(getattr(item[1], "left", 0)),
                item[0],
            ),
        )
    ]


def _shape_priority(shape: Any) -> int:
    if _is_title_shape(shape):
        return 0
    if bool(getattr(shape, "is_placeholder", False)):
        return 1
    return 2


def _is_title_shape(shape: Any) -> bool:
    if not bool(getattr(shape, "is_placeholder", False)):
        return False
    try:
        placeholder_type = shape.placeholder_format.type
    except (AttributeError, ValueError):
        return False
    return _enum_name(placeholder_type) in {"TITLE", "CENTER_TITLE"}


def _text_frame_text(text_frame: Any) -> str:
    lines: list[str] = []
    for paragraph in text_frame.paragraphs:
        text = _clean_text(paragraph.text)
        if not text:
            continue
        level = max(_int_value(getattr(paragraph, "level", 0)), 0)
        prefix = "- " if _paragraph_is_bullet(paragraph) else ""
        lines.append(f"{'  ' * level}{prefix}{text}")
    return "\n".join(lines)


def _paragraph_is_bullet(paragraph: Any) -> bool:
    paragraph_properties = getattr(paragraph._p, "pPr", None)
    if paragraph_properties is None:
        return False
    return any(
        str(child.tag).rsplit("}", 1)[-1] in {"buAutoNum", "buBlip", "buChar"}
        for child in paragraph_properties
    )


def _clean_text(value: Any) -> str:
    return "\n".join(
        line.rstrip() for line in str(value or "").replace("\r", "").split("\n")
    ).strip()


def _append_block(
    state: _PptxState,
    *,
    slide_index: int,
    block_type: str,
    text: str,
    shape: Any,
    page_box: list[int],
) -> str:
    block_index = sum(
        block.get("page") == slide_index for block in state.source_blocks
    )
    component_id = f"/page/{slide_index}/{block_type}/{block_index}"
    state.source_blocks.append(
        {
            "component_id": component_id,
            "page": slide_index,
            "block_index": block_index,
            "type": block_type,
            "text": text,
            "source": AXIOM_NATIVE_BLOCK_SOURCE,
            "parser_source": "pptx_native",
            "bbox": _shape_box(shape),
            "page_bbox": list(page_box),
        }
    )
    return component_id


def _shape_box(shape: Any) -> list[int]:
    left = _int_value(getattr(shape, "left", 0))
    top = _int_value(getattr(shape, "top", 0))
    width = _int_value(getattr(shape, "width", 0))
    height = _int_value(getattr(shape, "height", 0))
    return [left, top, left + width, top + height]


def _picture_description(shape: Any) -> str:
    candidates: list[str] = []
    for element in shape._element.iter():
        if str(element.tag).rsplit("}", 1)[-1] != "cNvPr":
            continue
        candidates.extend(
            str(element.get(field) or "").strip()
            for field in ("descr", "title")
        )
    candidates.append(str(getattr(shape, "name", "") or "").strip())
    for candidate in candidates:
        if candidate and not _GENERIC_PICTURE_NAME.fullmatch(candidate):
            return candidate
    return ""


def _markdown_table(headers: list[str], rows: list[list[str]]) -> str:
    def render(row: list[str]) -> str:
        return "| " + " | ".join(_markdown_cell(cell) for cell in row) + " |"

    return "\n".join(
        [
            render(headers),
            render(["---"] * len(headers)),
            *(render(row) for row in rows),
        ]
    )


def _markdown_cell(value: Any) -> str:
    return str(value or "").replace("\\", "\\\\").replace("|", "\\|").replace(
        "\n", "<br>"
    )


def _title_citations(state: _PptxState) -> list[str]:
    if state.first_title is None:
        return []
    return [
        str(block["component_id"])
        for block in state.source_blocks
        if block.get("type") == "Title"
    ][:1]


def _safe_extension(value: Any) -> str:
    extension = str(value or "bin").strip().lower().lstrip(".")
    return extension if re.fullmatch(r"[a-z0-9]+", extension) else "bin"


def _enum_name(value: Any) -> str:
    name = getattr(value, "name", None)
    if name:
        return str(name).upper()
    return str(value).split(" ", 1)[0].split("(", 1)[0].upper()


def _int_value(value: Any) -> int:
    try:
        return int(value or 0)
    except (TypeError, ValueError):
        return 0


def _optional_str(value: Any) -> str | None:
    if value is None:
        return None
    text = str(value).strip()
    return text or None
