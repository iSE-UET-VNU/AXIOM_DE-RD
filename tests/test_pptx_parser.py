from __future__ import annotations

import importlib.util
import json
import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch

from src.artifacts import write_ingested_artifacts
from src.ingestion.parsing.backends import PptxConfig, PptxParserBackend
from src.ingestion.parsing.chandra2 import Chandra2Provider
from src.ingestion.runner import run as run_ingestion
from src.models import DataObject, ParseStatus, PipelineState


PPTX_AVAILABLE = importlib.util.find_spec("pptx") is not None
PIL_AVAILABLE = importlib.util.find_spec("PIL") is not None


def _data_object(path: Path) -> DataObject:
    return DataObject(
        object_id="pptx-document",
        uri=path.as_posix(),
        content_type=(
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ),
        metadata={"format": "pptx", "file_name": path.name},
    )


@unittest.skipUnless(PPTX_AVAILABLE and PIL_AVAILABLE, "PPTX test dependencies missing")
class PptxParserTests(unittest.TestCase):
    def _presentation(self, root: Path) -> Path:
        from PIL import Image
        from pptx import Presentation
        from pptx.util import Inches

        image_path = root / "architecture.png"
        Image.new("RGB", (320, 180), color=(40, 100, 180)).save(image_path)

        presentation = Presentation()
        title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        title_slide.shapes.title.text = "Kiến trúc hệ thống"
        title_slide.placeholders[1].text = "Native PPTX parser"

        content_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        textbox = content_slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(4.0), Inches(1.2)
        )
        textbox.text_frame.text = "Mục tiêu"
        paragraph = textbox.text_frame.add_paragraph()
        paragraph.text = "Giữ nguyên Unicode"
        paragraph.level = 1

        table_shape = content_slide.shapes.add_table(
            2, 2, Inches(0.5), Inches(2.0), Inches(4.0), Inches(1.5)
        )
        table_shape.table.cell(0, 0).text = "Mô hình"
        table_shape.table.cell(0, 1).text = "Điểm"
        table_shape.table.cell(1, 0).text = "A"
        table_shape.table.cell(1, 1).text = "91%"

        first_picture = content_slide.shapes.add_picture(
            str(image_path), Inches(5.0), Inches(1.0), width=Inches(2.0)
        )
        first_picture._element.nvPicPr.cNvPr.set("descr", "Sơ đồ kiến trúc")
        content_slide.shapes.add_picture(
            str(image_path), Inches(7.2), Inches(1.0), width=Inches(2.0)
        )

        group = content_slide.shapes.add_group_shape()
        group.shapes.add_textbox(
            Inches(1.0), Inches(4.5), Inches(3.0), Inches(0.6)
        ).text = "Nội dung trong group"
        presentation.slides.add_slide(presentation.slide_layouts[6])

        path = root / "sample.pptx"
        presentation.save(path)
        return path

    def test_extracts_canonical_text_tables_figures_and_blocks(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            path = self._presentation(root)
            parser = PptxParserBackend(
                PptxConfig(output_dir=str(root / "assets"), project_root=str(root))
            )

            first = parser.parse(path, _data_object(path))
            second = parser.parse(path, _data_object(path))

        self.assertEqual(first.status, ParseStatus.SUCCESS)
        parsed = first.parsed_data
        self.assertIsNotNone(parsed)
        extraction = parsed.rows[0]["extraction"]
        self.assertEqual(extraction["document_type"], "presentation")
        self.assertEqual(extraction["title"], "Kiến trúc hệ thống")
        self.assertIn("Native PPTX parser", extraction["main_text"])
        self.assertIn("Nội dung trong group", extraction["main_text"])
        self.assertEqual(len(extraction["tables"]), 1)
        self.assertEqual(extraction["tables"][0]["headers"], ["Mô hình", "Điểm"])
        self.assertEqual(extraction["tables"][0]["rows"], [["A", "91%"]])
        self.assertEqual(parsed.tables[0].rows, [["A", "91%"]])
        self.assertEqual(len(extraction["figures"]), 2)
        self.assertEqual(extraction["figures"][0]["description"], "Sơ đồ kiến trúc")
        self.assertEqual(parsed.metadata["page_count"], 3)
        self.assertEqual(len(parsed.metadata["image_files"]), 2)
        self.assertEqual(
            parsed.metadata["image_files"][0]["path"],
            parsed.metadata["image_files"][1]["path"],
        )
        self.assertNotEqual(
            parsed.metadata["image_files"][0]["source_ref"],
            parsed.metadata["image_files"][1]["source_ref"],
        )
        self.assertEqual(
            parsed.rows[0]["reading_order"],
            second.parsed_data.rows[0]["reading_order"],
        )
        self.assertEqual(
            [item["name"] for item in parsed.metadata["image_files"]],
            [item["name"] for item in second.parsed_data.metadata["image_files"]],
        )

    def test_ingestion_filters_images_and_writes_common_artifact(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            path = self._presentation(root)
            assets = root / "ingested" / "assets"
            with patch.object(
                Chandra2Provider,
                "parse_file",
                side_effect=AssertionError("Chandra2 must not receive PPTX"),
            ):
                output = run_ingestion(
                    path,
                    parser_config={
                        "provider": "chandra2",
                        "pptx": {
                            "output_dir": str(assets),
                            "project_root": str(root),
                        },
                    },
                    project_root=root,
                )

            parsed = output.parsed_data[0]
            self.assertEqual(parsed.metadata["backend"], "pptx")
            self.assertEqual(parsed.metadata["image_filtering"]["kept_count"], 2)
            self.assertEqual(parsed.metadata["image_filtering"]["copied_count"], 2)
            self.assertTrue(
                all(
                    Path(root / item["path"]).is_file()
                    and Path(item["path"]).parent.name == "filtered_images"
                    for item in parsed.metadata["image_files"]
                )
            )

            state = PipelineState(
                run_id="pptx-test",
                input_source=path.as_posix(),
                embedded_dir="embedded",
                output_dir="output",
                data_objects=output.data_objects,
                parsed_data=output.parsed_data,
                initial_schemas=output.initial_schemas,
                ingestion_config={"provider": "chandra2"},
            )
            artifact_dir = root / "ingested"
            write_ingested_artifacts(state, artifact_dir, project_root=root)
            document_path = next((artifact_dir / "documents").glob("*.json"))
            payload = json.loads(document_path.read_text(encoding="utf-8"))

        self.assertEqual(payload["contract_version"], "ingested-document-v2")
        extraction = payload["parsed"]["rows"][0]["extraction"]
        self.assertIn("Kiến trúc hệ thống", extraction["main_text"])
        self.assertEqual(len(extraction["tables"]), 1)
        self.assertEqual(len(extraction["figures"]), 2)

    def test_corrupt_pptx_returns_failed_result(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            path = Path(temp_dir) / "broken.pptx"
            path.write_bytes(b"not an office package")
            result = PptxParserBackend().parse(path, _data_object(path))

        self.assertEqual(result.status, ParseStatus.FAILED)
        self.assertEqual(result.backend, "pptx")
        self.assertTrue(result.error.get("message"))


if __name__ == "__main__":
    unittest.main()
